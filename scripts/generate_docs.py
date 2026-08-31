import os
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    run = h.runs[0]
    run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB) # Brand blue

def generate_word_doc():
    doc = Document()
    
    # 1. Title Page
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("ARO Cluster Upgrade Automation\n")
    run.font.size = Pt(28)
    run.bold = True
    run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    
    sub = doc.add_paragraph("Rule-Based Ansible Automation for ARO Y-Stream Upgrades")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub.runs[0]
    run.font.size = Pt(16)
    
    doc.add_paragraph("\nAuthor: Technical Writer / Solutions Architect")
    doc.add_paragraph("Date: August 2026")
    doc.add_paragraph("Version: 1.0\n")
    
    # Legend
    p = doc.add_paragraph()
    p.add_run("Status Legend:\n").bold = True
    p.add_run("PASS (Green)").font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)
    p.add_run(" | ")
    p.add_run("WARN (Amber)").font.color.rgb = RGBColor(0xD9, 0x77, 0x06)
    p.add_run(" | ")
    p.add_run("FAIL/HARD (Red)").font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    
    doc.add_page_break()
    
    # 2. Executive Summary
    add_heading(doc, "Executive Summary", 1)
    doc.add_paragraph("The ARO Cluster Upgrade Automation suite is a rule-based, deterministic Ansible automation solution designed to perform Y-stream (minor) upgrades of Azure Red Hat OpenShift clusters as ordered, sequential version hops (e.g. 4.18.09 → 4.19.15 → 4.20.08).")
    doc.add_paragraph("It addresses the risk, inconsistency, and lack of auditability inherent in manual upgrades. Key benefits include:")
    doc.add_paragraph("- Deterministic: Rule-based, no AI in the execution path. All decisions are hardcoded rules.", style='List Bullet')
    doc.add_paragraph("- Auditable: Produces dual .txt/.csv logs and colour-coded client HTML reports.", style='List Bullet')
    doc.add_paragraph("- Dual-Version Safe: Runs unchanged on both Ansible 2.7.17 (test) and 2.14.18 (prod).", style='List Bullet')
    doc.add_paragraph("- Safe-Fail: Login once, logout on any failure (via block/rescue). HARD gates halt execution.", style='List Bullet')
    
    # 3. Table of Contents
    add_heading(doc, "Table of Contents", 1)
    doc.add_paragraph("1. Executive Summary\n2. Solution Overview\n3. Architecture\n4. End-to-End Workflow\n5. Phase Deep-Dives\n6. Validation Model\n7. Upgrade Hop Mechanics\n8. Monitoring & Notifications\n9. Operator Upgrade & Validation\n10. Reporting, Logging & Outputs\n11. Dual-Version Compatibility & Migration\n12. Build Units Reference\n13. Appendix")
    doc.add_page_break()
    
    # 4. Solution Overview
    add_heading(doc, "Solution Overview", 1)
    doc.add_paragraph("Goals: Automate multi-version Y-stream ARO upgrades deterministically without skipping minor versions. Enforce 13-check prevalidation, live monitoring, and 10-check postvalidation.")
    doc.add_paragraph("Core User Flow:")
    doc.add_paragraph("1. Define the upgrade path in vars.\n2. Run CLI entrypoint (00_Run.sh).\n3. Policy check & baseline.\n4. Prevalidation hard gate.\n5. Execute hops.\n6. Monitor each hop.\n7. Settle-gate between hops.\n8. Postvalidation.\n9. Deliver & log out.")
    doc.add_paragraph("In-Scope: Y-stream sequential hops, 13-check prevalidation, live monitoring, HTML reports, logs, snapshot.")
    doc.add_paragraph("Out-of-Scope: AI, GitOps/Argo, cluster backup/restore, Z-stream/X-stream upgrades.")
    
    # 5. Architecture
    add_heading(doc, "Architecture", 1)
    doc.add_paragraph("System structure relies on Ansible for orchestration, Bash for the CLI entrypoint, oc for cluster interfacing, jq for data parsing, and sendmail for notifications. All paths are dynamic and secrets are variable references.")
    try:
        doc.add_picture("DIAGRAMS/architecture-diagram.png", width=Inches(6.0))
    except Exception as e:
        doc.add_paragraph(f"[Architecture Diagram Missing: {e}]")
        
    doc.add_paragraph("Invariants:\n- No AI in execution path.\n- Codebase runs unchanged on Ansible 2.7.17 and 2.14.18.\n- Login once, logout on any failure.\n- oc live queries only; no cross-phase state persistence (except Phase 01 snapshot).\n- Minor versions are never skipped.\n- HARD gates halt the chain.")
    
    # 6. End-to-End Workflow
    add_heading(doc, "End-to-End Workflow", 1)
    doc.add_paragraph("The master workflow traverses the following lifecycle: Login → Snapshot → Secrets → Prevalidation Gate → Hop Loop → Postvalidation → Operator Upgrade → Logout.")
    try:
        doc.add_picture("DIAGRAMS/master-workflow.png", width=Inches(6.0))
    except Exception:
        doc.add_paragraph("[Master Workflow Diagram Missing]")
        
    # 7. Phase Deep-Dives
    add_heading(doc, "Phase Deep-Dives", 1)
    doc.add_heading("Phase 01: Policy Check & Baseline", 2)
    doc.add_paragraph("Purpose: Login once, capture baseline JSON snapshot, validate policy/edge.\nInputs: Secrets, upgrade_path.\nActions: oc login, oc get clusterversion/operators/nodes/routes -o json.\nGate: HARD if invalid edge.")
    doc.add_heading("Phase 02: Prevalidation", 2)
    doc.add_paragraph("Purpose: 13-check hard gate before any upgrade.\nActions: oc/jq queries for CO, Node, MCP, etcd, etc.\nGate: HARD halts chain; WARN surfaces in report.")
    doc.add_heading("Phase 03 & 04: Upgrade Hops & Live Monitoring", 2)
    doc.add_paragraph("Purpose: Trigger and monitor each hop.\nActions: oc adm upgrade --to, monitor cv/mcp/nodes every 2m.")
    doc.add_heading("Phase 05: Postvalidation", 2)
    doc.add_paragraph("Purpose: 10-check postvalidation and baseline diff.\nGate: HARD.")
    doc.add_heading("Unit 20: Operator Upgrade", 2)
    doc.add_paragraph("Purpose: Check OLM compatibility, approve InstallPlans, wait for Succeeded.")
    
    # 8. Validation Model
    add_heading(doc, "Validation Model", 1)
    doc.add_paragraph("13 Prevalidation Checks:")
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '#'
    hdr_cells[1].text = 'Check'
    hdr_cells[2].text = 'Rule'
    hdr_cells[3].text = 'Gate'
    checks = [
        ("01", "ClusterOperators Health", "Available=True, Degraded=False", "HARD"),
        ("02", "Node Health", "Ready=True, Schedulable", "HARD"),
        ("03", "MachineConfigPool Health", "Updated=True, Degraded=False", "HARD"),
        ("04", "etcd Member Health", "Running/Ready, 0 degraded", "HARD"),
        ("05", "Admin Acknowledgements", "Upgradeable condition / acks", "HARD"),
        ("06", "Target Upgrade Edge", "Valid available/conditional edge", "HARD"),
        ("07", "Resource Utilization", "CPU/Mem < 90%", "HARD"),
        ("08", "Pending CSRs", "No unapproved CSRs", "WARN"),
        ("09", "Disk Pressure", "No DiskPressure=True", "WARN"),
        ("10", "PDB Risk", "Drain headroom available", "WARN"),
        ("11", "Node Drain Headroom", "Utilization if largest node evicted", "WARN"),
        ("12", "Critical Pods", "openshift-* pods Running", "WARN"),
        ("13", "Firing Alerts", "No critical Prometheus alerts", "WARN")
    ]
    for n, c, r, g in checks:
        row = table.add_row().cells
        row[0].text = n
        row[1].text = c
        row[2].text = r
        row[3].text = g
        if g == "HARD":
            row[3].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        else:
            row[3].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xD9, 0x77, 0x06)
            
    doc.add_paragraph("\n10 Postvalidation Checks run after all hops, including a baseline diff.")
    
    # 9. Upgrade Hop Mechanics
    add_heading(doc, "Upgrade Hop Mechanics", 1)
    doc.add_paragraph("Each hop executes natively via oc adm upgrade --to. Includes setting the channel, verifying the edge, applying admin-ack, triggering upgrade, monitoring, and passing the settle-gate (CV at target, CO Available, MCP Updated).")
    try:
        doc.add_picture("DIAGRAMS/per-hop-upgrade.png", width=Inches(6.0))
    except Exception:
        doc.add_paragraph("[Per-Hop Diagram Missing]")
        
    # 10. Monitoring & Notifications
    add_heading(doc, "Monitoring & Notifications", 1)
    doc.add_paragraph("Polling cadence is 2 minutes. HTML heartbeat emails are sent every 20 minutes, with immediate change alerts on state mutation. Timeout guard is typically 90 minutes per hop.")
    try:
        doc.add_picture("DIAGRAMS/monitoring-loop.png", width=Inches(6.0))
    except Exception:
        doc.add_paragraph("[Monitoring Diagram Missing]")
        
    # 11. Operator Upgrade & Validation (Unit 20)
    add_heading(doc, "Operator Upgrade & Validation (Unit 20)", 1)
    doc.add_paragraph("Compatibility scan → InstallPlan approval → CSV settle → final validation.")
    try:
        doc.add_picture("DIAGRAMS/operator-upgrade.png", width=Inches(6.0))
    except Exception:
        doc.add_paragraph("[Operator Upgrade Diagram Missing]")
        
    # 12. Reporting, Logging & Outputs
    add_heading(doc, "Reporting, Logging & Outputs", 1)
    doc.add_paragraph("Reports: HTML prevalidation and postvalidation reports in output/.\nLogs: Dual .txt and .csv run logs in logs/.\nSnapshot: Baseline JSON snapshot in snapshots/.")
    
    # 13. Dual-Version Compatibility
    add_heading(doc, "Dual-Version Compatibility & Migration", 1)
    doc.add_paragraph("Codebase runs on Ansible 2.7.17 and 2.14.18. Utilizes short module names, include_tasks, | default(''), and a dual-delivery sendmail contract.")
    
    # 14. Build Units Reference
    add_heading(doc, "Build Units Reference", 1)
    try:
        doc.add_picture("DIAGRAMS/build-units-roadmap.png", width=Inches(6.0))
    except Exception:
        doc.add_paragraph("[Roadmap Diagram Missing]")
        
    # 15. Appendix
    add_heading(doc, "Appendix", 1)
    doc.add_paragraph("Glossary:\nARO: Azure Red Hat OpenShift\nMCP: MachineConfigPool\nCO: ClusterOperator\nCSV: ClusterServiceVersion\nInstallPlan: OLM upgrade instruction\nEdge: Supported upgrade path between versions\nY-stream: Minor version release")
    
    doc.save("ARO-Upgrade-Automation-Documentation.docx")

def add_slide(prs, title, content_lines, img_path=None):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, PptxInches(1.5), PptxInches(1.5), width=PptxInches(10))
    else:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(content_lines):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            
    # Add generic speaker note
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = f"Speaker Notes for {title}:\nDiscuss the key points shown on the slide. Emphasize rule-based determinism and safe-fail mechanisms."
    return slide

def generate_powerpoint():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ARO Cluster Upgrade Automation"
    subtitle.text = "Rule-Based Ansible Automation for ARO Y-Stream Upgrades\nBy: Technical Writer\nAugust 2026"
    
    # 2. The Problem
    add_slide(prs, "The Problem", ["Manual ARO upgrades are risky and inconsistent.", "Lack of audit trails.", "High potential for human error."])
    
    # 3. The Solution
    add_slide(prs, "The Solution", ["Deterministic, rule-based Ansible automation.", "Zero AI in the execution path.", "Safe, ordered sequential version hops."])
    
    # 4. Key Benefits
    add_slide(prs, "Key Benefits", ["Deterministic execution.", "Fully auditable (.txt, .csv, HTML).", "Dual-version compatibility (2.7.17 / 2.14.18).", "Safe-fail: Login once, logout on any failure."])
    
    # 5. Architecture at a Glance
    add_slide(prs, "Architecture at a Glance", [], img_path="DIAGRAMS/architecture-diagram.png")
    
    # 6. End-to-End Workflow
    add_slide(prs, "End-to-End Workflow", [], img_path="DIAGRAMS/master-workflow.png")
    
    # 7. Phase Map
    add_slide(prs, "Phase Map", ["Phase 01: Login & Baseline Snapshot", "Phase 02: Prevalidation", "Phase 03: Initiate Upgrade Loop", "Phase 04: Live Monitoring", "Phase 05: Postvalidation", "Unit 20: Operator Upgrade"])
    
    # 8. Prevalidation Hard Gate
    add_slide(prs, "Prevalidation Hard Gate", [], img_path="DIAGRAMS/prevalidation-gate.png")
    
    # 9. Upgrade Hops
    add_slide(prs, "Upgrade Hops", [], img_path="DIAGRAMS/per-hop-upgrade.png")
    
    # 10. Monitoring & Alerts
    add_slide(prs, "Monitoring & Alerts", [], img_path="DIAGRAMS/monitoring-loop.png")
    
    # 11. Postvalidation & Baseline Diff
    add_slide(prs, "Postvalidation & Baseline Diff", ["10 Checks executed after final hop.", "Ensures CV, CO, MCP, Node, and etcd health.", "Compares Phase 05 state against Phase 01 baseline snapshot."])
    
    # 12. Operator Upgrade & Validation
    add_slide(prs, "Operator Upgrade & Validation", [], img_path="DIAGRAMS/operator-upgrade.png")
    
    # 13. Safety & Guardrails
    add_slide(prs, "Safety & Guardrails", ["Login-once, logout-on-failure.", "HARD gates immediately halt execution.", "WARN gates notify but proceed.", "force flag is manual-only."])
    
    # 14. Dual-Version Strategy
    add_slide(prs, "Dual-Version Strategy", ["Runs on Ansible 2.7.17 and 2.14.18.", "Short module names.", "No bare includes.", "Jinja2 strictness via default('')."])
    
    # 15. Reporting & Outputs
    add_slide(prs, "Reporting & Outputs", ["HTML Reports (Pre/Post validation).", "Logs (dual .txt and .csv).", "Baseline JSON snapshot."])
    
    # 16. Build Roadmap
    add_slide(prs, "Build Roadmap", [], img_path="DIAGRAMS/build-units-roadmap.png")
    
    # 17. Summary & Next Steps
    add_slide(prs, "Summary & Next Steps", ["Testing on RHEL 8 jump host.", "End-to-end dry-run validation.", "Operator handover."])
    
    # 18. Q&A
    add_slide(prs, "Q&A / Thank You", ["Any Questions?"])
    
    prs.save("ARO-Upgrade-Automation-Deck.pptx")

if __name__ == "__main__":
    generate_word_doc()
    generate_powerpoint()
    print("Generated ARO-Upgrade-Automation-Documentation.docx and ARO-Upgrade-Automation-Deck.pptx successfully.")
